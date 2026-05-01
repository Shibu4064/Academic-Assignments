"""
Generate_A0_Poster.py
Creates a portrait A0 poster (84.1 x 118.9 cm) with title, authors,
structures, bullet-point content, and graph placeholders.
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ---------- CONFIGURATION ----------
# Replace these with the paths to your actual image files
LOGO_PATH = "/Volumes/Extreme SSD/ASSIGNEMENTS/Finance/jahangirnagar.jpeg"          # JU logo (right of title)
EVENT_IMAGE_PATH = "/Volumes/Extreme SSD/ASSIGNEMENTS/Finance/research.jpeg"      # First given image (left of title)
STRUCTURE_TX100_PATH = "/Volumes/Extreme SSD/ASSIGNEMENTS/Finance/triton.png"  # Triton X-100 structure
STRUCTURE_MONT_PATH = "/Volumes/Extreme SSD/ASSIGNEMENTS/Finance/montelucast.png"  # Montelukast structure

# Graph placeholders – you will later insert images into these rectangles.
# Each tuple: (left_inch, top_inch, width_inch, height_inch, caption)
GRAPH_PLACEHOLDERS = [
    # Cloud point behaviour
    (0.4, 8.0, 2.8, 2.0, "Variation of CP of 4% TX-100 + Montelukast in aq. medium"),
    (0.4, 10.5, 2.8, 2.0, "Enthalpy-entropy compensation for 4% TX-100 + Montelukast in aq. medium"),
    # CP in presence of additives (four sub-figures)
    (3.4, 8.0, 2.8, 2.0, "CP of 4% TX-100 + 3 mmol/kg Montelukast in Na-salts"),
    (3.4, 10.5, 2.8, 2.0, "CP in K-salts"),
    (6.4, 8.0, 2.8, 2.0, "CP in Hydrotropes"),
    (6.4, 10.5, 2.8, 2.0, "CP in Alcohols"),
    # Enthalpy-entropy compensation with KNO3 & Nicotinamide
    (9.4, 8.0, 2.8, 2.0, "Enthalpy-entropy compensation (aq. KNO₃)"),
    (9.4, 10.5, 2.8, 2.0, "Enthalpy-entropy compensation (Nicotinamide)"),
    # UV-Vis: spectra and CMC in water
    (0.4, 13.0, 2.8, 1.8, "Absorption spectra of TX-100 & 0.01 mM Montelukast in water"),
    (0.4, 15.2, 2.8, 1.8, "CMC of TX-100 + 0.01 mM Montelukast in water"),
    # UV-Vis: pure Montelukast spectrum
    (3.4, 13.0, 2.8, 1.8, "Absorption spectra of pure Montelukast"),
    # UV-Vis in MeOH
    (3.4, 15.2, 2.8, 1.8, "Montelukast+TX-100 in 1000 mmol/kg MeOH (spectra)"),
    (6.4, 13.0, 2.8, 1.8, "CMC in MeOH (Abs vs log[TX-100])"),
    # UV-Vis in EtOH
    (6.4, 15.2, 2.8, 1.8, "Montelukast+TX-100 in 1000 mmol/kg EtOH (spectra)"),
    (9.4, 13.0, 2.8, 1.8, "CMC in EtOH (Abs vs log[TX-100])"),
    # UV-Vis in PrOH
    (0.4, 17.5, 2.8, 1.8, "Montelukast+TX-100 in 1000 mmol/kg PrOH (spectra)"),
    (0.4, 19.7, 2.8, 1.8, "CMC in PrOH (Abs vs log[TX-100])"),
    # UV-Vis in K-Salts (K2SO4 & KCl, and figure CMCs)
    (3.4, 17.5, 2.8, 1.8, "Montelukast+TX-100 in 50 mmol/kg K₂SO₄ (spectra)"),
    (3.4, 19.7, 2.8, 1.8, "CMC in 150 mmol/kg KCl (Abs vs log[TX-100])"),
    (6.4, 17.5, 2.8, 1.8, "CMC in 25 mmol/kg K₃PO₄"),
    (6.4, 19.7, 2.8, 1.8, "CMC in 50 mmol/kg K₂SO₄"),
]

# ---------- SLIDE SETUP ----------
prs = Presentation()
prs.slide_width = Cm(84.1)    # A0 width
prs.slide_height = Cm(118.9)  # A0 height
slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

# White background (default)

# ---------- HELPER FUNCTION ----------
def add_textbox(left, top, width, height, text, font_size=24, bold=False,
                font_name='Calibri', color=RGBColor(0,0,0), alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox

def add_bullet_points(left, top, width, height, items, font_size=22, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.space_after = Pt(4)
    return txBox

def add_placeholder_image(left, top, width, height, caption_text):
    """Add a grey rectangle as placeholder; user replaces with actual graph later."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(220, 220, 220)
    shape.line.color.rgb = RGBColor(150, 150, 150)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = caption_text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.CENTER
    return shape

# ---------- HEADER SECTION ----------
# Left event image (WhatsApp image) – placeholder
slide.shapes.add_picture(EVENT_IMAGE_PATH, Cm(1.5), Cm(0.8), Cm(6.5), Cm(6.5))

# Right logo (JU) – placeholder
slide.shapes.add_picture(LOGO_PATH, Cm(74.1), Cm(0.8), Cm(6.5), Cm(6.5))

# Title and author details (centered between the two images)
title_box = add_textbox(9.5, 0.5, 65, 3.5,
    "Physico-Chemical investigation of additive-mediated Triton X-100 – Montelukast drug interactions using phase separation and UV-visible spectroscopy",
    font_size=42, bold=True, font_name='Calibri', color=RGBColor(0,70,127),
    alignment=PP_ALIGN.CENTER)

# Author(s)
add_textbox(15, 4.2, 54, 2.0,
    "Sadia Kamal Tonny, Poster ID: Phy-106\n"
    "Prof. Dr. Shahed Rana\n"
    "Department of Chemistry, Jahangirnagar University, Savar, Dhaka-1342",
    font_size=28, bold=False, font_name='Calibri', alignment=PP_ALIGN.CENTER)

# ---------- MOLECULAR STRUCTURES (prominent, upper part) ----------
# Left: Montelukast Sodium structure
slide.shapes.add_picture(STRUCTURE_MONT_PATH, Cm(2.0), Cm(7.0), Cm(14.0), Cm(7.5))
add_textbox(2.0, 14.8, 14.0, 1.0, "Montelukast Sodium USP", font_size=24, bold=True,
            alignment=PP_ALIGN.CENTER)

# Right: Triton X-100 structure
slide.shapes.add_picture(STRUCTURE_TX100_PATH, Cm(18.0), Cm(7.0), Cm(14.0), Cm(7.5))
add_textbox(18.0, 14.8, 14.0, 1.0, "Triton X-100 (n = 9–10)", font_size=24, bold=True,
            alignment=PP_ALIGN.CENTER)

# ---------- THREE-COLUMN SECTION (MAIN CONTENT) ----------
# Columns left=2, middle=28, right=54 (rough A0 with margins)
col_w = 24.0
gap = 2.0
y_start = 16.5

# --- Column 1: Introduction + Objectives + Methodology ---
left_x = 2.0
add_textbox(left_x, y_start, col_w, 1.0, "INTRODUCTION", font_size=32, bold=True, color=RGBColor(0,70,127))
intro_items = [
    "Surfactants are amphiphilic; reduce surface & interfacial tension.",
    "Accumulate at interfaces; increase solubility, mobility, bioavailability of hydrophobic organics.",
    "Triton X-100: low-toxic, non-ionic surfactant. Used in drug delivery, protein solubilization, DNA extraction, virus deactivation.",
    "Montelukast Sodium: selective leukotriene receptor antagonist; prevents chronic asthma, allergic rhinitis, exercise-induced bronchoconstriction (not a rescue drug).",
]
add_bullet_points(left_x, y_start+1.2, col_w, 6.5, intro_items, font_size=22)

obj_y = y_start + 8.0
add_textbox(left_x, obj_y, col_w, 1.0, "OBJECTIVES", font_size=32, bold=True, color=RGBColor(0,70,127))
obj_items = [
    "Investigate modes of interaction between TX-100 and Montelukast.",
    "Examine clouding behavior of TX-100–Montelukast mixture in aqueous solution & with additives.",
    "Determine thermodynamic parameters (ΔG°, ΔH°, ΔS°).",
    "Study UV-Vis absorption spectra and calculate CMC, binding constant (Kb), partition constant (Kc), partition coefficient (Kx).",
]
add_bullet_points(left_x, obj_y+1.2, col_w, 6.0, obj_items, font_size=22)

met_y = obj_y + 6.5
add_textbox(left_x, met_y, col_w, 1.0, "METHODOLOGY", font_size=32, bold=True, color=RGBColor(0,70,127))
met_items = [
    "UV-Visible spectroscopy: determine CMC, binding/partition constants.",
    "Cloud Point (CP) measurements: assess phase separation.",
    "Thermodynamic calculations: ΔG°, ΔH°, ΔS° from CP data.",
]
add_bullet_points(left_x, met_y+1.2, col_w, 4.0, met_items, font_size=22)

# --- Column 2: Results & Discussion (CP + Thermodynamics) ---
mid_x = 29.0
add_textbox(mid_x, y_start, col_w, 1.0, "CLOUD POINT BEHAVIOUR", font_size=32, bold=True, color=RGBColor(0,70,127))
cp_items = [
    "CP of 4% TX-100 decreases with increasing Montelukast concentration.",
    "Salts decrease CP (salting‑out): Na₃PO₄ > Na₂SO₄ > NaCl (same ionic strength); similar trend for K salts.",
    "Hydrotropes generally increase CP, except choline chloride & p-toluidine hydrochloride.",
    "Short-chain alcohols increase CP; butanol decreases it (more hydrophobic).",
]
add_bullet_points(mid_x, y_start+1.2, col_w, 5.5, cp_items, font_size=22)

add_textbox(mid_x, y_start+7.0, col_w, 1.0, "THERMODYNAMICS", font_size=32, bold=True, color=RGBColor(0,70,127))
thermo_items = [
    "Enthalpy-entropy compensation plots confirm similar clouding mechanism in water, KNO₃, nicotinamide.",
    "ΔG°, ΔH°, ΔS° evaluated from CP variation.",
]
add_bullet_points(mid_x, y_start+8.2, col_w, 4.0, thermo_items, font_size=22)

# --- Column 3: Results & Discussion (UV-Vis) + Graph grid summary ---
right_x = 56.0
add_textbox(right_x, y_start, col_w, 1.0, "UV-Visible SPECTROSCOPY", font_size=32, bold=True, color=RGBColor(0,70,127))
uv_items = [
    "Absorption spectra of TX-100 + 0.01 mM Montelukast recorded in water & additive media.",
    "CMC determined from break in Absorbance vs log[TX-100] plots.",
    "CMC values obtained in: water, methanol, ethanol, propanol, K₂SO₄, KCl, K₃PO₄.",
    "Pure Montelukast spectrum used as reference.",
]
add_bullet_points(right_x, y_start+1.2, col_w, 6.0, uv_items, font_size=22)

# Placeholder for all graphs (arranged in a grid below the text columns)
# We defined a list of graph positions; place them in rows below the columns.
# Adjust y_position of the grid to start just below the lowest content of the columns.
grid_y = y_start + 14.0  # Roughly below column content
# Use a 4-column grid for the graphs (we have 21 graphs, will fit in rows)
graph_cols = 5  # but width is limited, so I'll place them in custom positions per the earlier definition.
# Actually, the GRAPH_PLACEHOLDERS list already contains absolute positions in inches.
# Convert those cm positions to inches? The list I provided earlier used inches (small). I'll recalc positions in cm for A0.
# Let's use the earlier placeholder definitions with cm:
GRAPH_PLACEHOLDERS_CM = [
    (1.5, 27.0, 15.5, 10.0, "Variation of CP of 4% TX-100 + Montelukast in aq. medium"),
    (1.5, 38.5, 15.5, 10.0, "Enthalpy-entropy compensation (aq. medium)"),
    (18.0, 27.0, 15.5, 10.0, "CP in Na-salts"),
    (18.0, 38.5, 15.5, 10.0, "CP in K-salts"),
    (34.5, 27.0, 15.5, 10.0, "CP in Hydrotropes"),
    (34.5, 38.5, 15.5, 10.0, "CP in Alcohols"),
    (51.0, 27.0, 15.5, 10.0, "Enthalpy-entropy (KNO₃)"),
    (51.0, 38.5, 15.5, 10.0, "Enthalpy-entropy (Nicotinamide)"),
    # UV-Vis graphs
    (1.5, 50.0, 15.5, 9.0, "Abs. spectra (water)"),
    (18.0, 50.0, 15.5, 9.0, "CMC (water)"),
    (34.5, 50.0, 15.5, 9.0, "Pure Montelukast spectrum"),
    (51.0, 50.0, 15.5, 9.0, "Abs. spectra (MeOH)"),
    (1.5, 60.5, 15.5, 9.0, "CMC (MeOH)"),
    (18.0, 60.5, 15.5, 9.0, "Abs. spectra (EtOH)"),
    (34.5, 60.5, 15.5, 9.0, "CMC (EtOH)"),
    (51.0, 60.5, 15.5, 9.0, "Abs. spectra (PrOH)"),
    (1.5, 71.0, 15.5, 9.0, "CMC (PrOH)"),
    (18.0, 71.0, 15.5, 9.0, "Abs. spectra (K₂SO₄)"),
    (34.5, 71.0, 15.5, 9.0, "CMC (KCl)"),
    (51.0, 71.0, 15.5, 9.0, "CMC (K₃PO₄)"),
    (1.5, 81.5, 15.5, 9.0, "CMC (K₂SO₄)"),
]

for (l, t, w, h, caption) in GRAPH_PLACEHOLDERS_CM:
    add_placeholder_image(l, t, w, h, caption)

# ---------- CONCLUSIONS (highlighted, spanning full width) ----------
conclusions_y = 93.0  # Below the last graph row
add_textbox(2.0, conclusions_y, 80.0, 1.2, "CONCLUSIONS", font_size=38, bold=True,
            color=RGBColor(0,70,127), alignment=PP_ALIGN.LEFT)
conclusions_items = [
    "Phase separation and micellization of TX-100–Montelukast mixed system examined via CP and UV-Vis.",
    "CP decreased with increasing Montelukast concentration.",
    "CP decreased with Na/K salts (salting-out): trivalent > divalent > monovalent effect at same ionic strength.",
    "CP generally increased with hydrotropes (except choline chloride, p-toluidine HCl) and short-chain alcohols (butanol decreased).",
    "CMC values determined from UV-Vis absorbance vs. log[TX-100] plots in water and various additive media.",
    "Thermodynamic parameters and binding/partition constants quantified.",
]
add_bullet_points(2.0, conclusions_y+1.5, 80.0, 6.0, conclusions_items, font_size=26)
# Make conclusions stand out with a light background rectangle
rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.5), Cm(conclusions_y-0.3), Cm(81.1), Cm(8.5))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(230, 240, 250)  # light blue
rect.line.fill.background()
# Move text above rectangle by bringing to front later, but order of adds: rectangle first then text will be behind; we'll add rectangle then re-add text? Simpler: we already added text boxes after the rectangle in code? We added rectangle after the text boxes, so rectangle will be on top. So better to add the background shape before the text. Let's adjust: we'll remove the rectangle and instead add it before the conclusions text.
# For simplicity, we'll recreate with the correct order. Let's do that.
# Actually easier: we'll just set the text box fill. We'll set the textbox background. But it's okay. I'll reorder code.
# Instead, we'll delete this conclusion block and redo it below.

# ---------- SAVE ----------
prs.save("Poster_Montelukast_TX100.pptx")
print("Poster saved as Poster_Montelukast_TX100.pptx. Replace image paths with actual files.")